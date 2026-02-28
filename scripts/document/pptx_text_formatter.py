#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文本格式自动修复工具 - PPTX版本
支持处理PowerPoint文档中的格式问题

功能列表：
1. 双引号格式修复
   将所有双引号统一为中文标准引号："和"
   奇数个引号 → " (中文左双引号 U+201C)
   偶数个引号 → " (中文右双引号 U+201D)

2. 英文标点符号转中文
   , → ，  (逗号)
   : → ：  (冒号)
   ; → ；  (分号)
   ! → ！  (感叹号)
   ? → ？  (问号)
   ( → （  (左括号)
   ) → ）  (右括号)

3. 中文单位转标准符号
   面积：平方米 → m²、平方公里 → km²
   体积：立方米 → m³、立方厘米 → cm³
   长度：公里 → km、厘米 → cm、毫米 → mm
   质量：公斤 → kg、毫克 → mg
   容量：毫升 → mL、微升 → μL
   时间：小时 → h、分钟 → min、秒钟 → s
   温度：摄氏度 → ℃、华氏度 → ℉

4. 单位上标格式化
   m2 → m²、m3 → m³、km2 → km²、km3 → km³

使用方法：
    python3 text_formatter.py 文件名.pptx
"""

import re
import sys
import shutil
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent.parent / "lib"))
from finder import get_input_files

try:
    from pptx import Presentation
except ImportError:
    print("❌ 错误: 缺少 python-pptx 库")
    print("💡 请运行: pip install python-pptx")
    sys.exit(1)


def fix_quotes(content):
    """
    替换所有双引号为中文标准引号
    奇数位置 → " (左引号)
    偶数位置 → " (右引号)
    """
    quote_pattern = r'["""''「」]'
    counter = 0
    
    def replace_quote(match):
        nonlocal counter
        counter += 1
        return '"' if counter % 2 == 1 else '"'
    
    result = re.sub(quote_pattern, replace_quote, content)
    return result, counter


def fix_punctuation(content):
    """
    将英文标点符号转换为中文标点符号
    """
    punctuation_map = {
        ',': '，',
        ':': '：',
        ';': '；',
        '!': '！',
        '?': '？',
        '(': '（',
        ')': '）',
    }
    
    result = content
    replacement_count = 0
    
    for eng_punct, chn_punct in punctuation_map.items():
        escaped_punct = re.escape(eng_punct)
        count = len(re.findall(escaped_punct, result))
        replacement_count += count
        result = re.sub(escaped_punct, chn_punct, result)
    
    return result, replacement_count


def fix_units(content):
    """
    将中文单位转换为标准符号单位
    """
    units_map = {
        # 面积单位
        '平方公里': 'km²',
        '平方千米': 'km²',
        '平方米': 'm²',
        '平方厘米': 'cm²',
        '平方毫米': 'mm²',
        # 体积单位
        '立方米': 'm³',
        '立方厘米': 'cm³',
        '立方毫米': 'mm³',
        '立方公里': 'km³',
        '立方千米': 'km³',
        # 长度单位
        '公里': 'km',
        '千米': 'km',
        '厘米': 'cm',
        '毫米': 'mm',
        '微米': 'μm',
        '纳米': 'nm',
        # 质量单位
        '公斤': 'kg',
        '千克': 'kg',
        '毫克': 'mg',
        '微克': 'μg',
        # 容量单位
        '毫升': 'mL',
        '微升': 'μL',
        # 时间单位
        '小时': 'h',
        '分钟': 'min',
        '秒钟': 's',
        # 温度单位
        '摄氏度': '℃',
        '华氏度': '℉',
        # m2/m3 转换
        'km2': 'km²',
        'km3': 'km³',
        'm2': 'm²',
        'm3': 'm³',
    }
    
    result = content
    replacement_count = 0
    
    sorted_units = sorted(units_map.items(), key=lambda x: len(x[0]), reverse=True)
    
    for chn_unit, symbol in sorted_units:
        count = result.count(chn_unit)
        if count > 0:
            replacement_count += count
            result = result.replace(chn_unit, symbol)
    
    return result, replacement_count


def process_text(text, stats):
    """
    处理文本，应用所有转换
    """
    if not text:
        return text
    
    result, quote_count = fix_quotes(text)
    result, punct_count = fix_punctuation(result)
    result, unit_count = fix_units(result)
    
    stats['quotes'] += quote_count
    stats['punctuation'] += punct_count
    stats['units'] += unit_count
    
    return result


def process_run(run, stats):
    """
    处理单个run的文本
    """
    if run.text:
        original = run.text
        fixed = process_text(original, stats)
        if fixed != original:
            run.text = fixed


def process_text_frame(text_frame, stats):
    """
    处理文本框中的所有段落和run
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            process_run(run, stats)


def process_table(table, stats):
    """
    处理表格中的所有单元格
    """
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                process_text_frame(cell.text_frame, stats)


def process_shape(shape, stats):
    """
    处理单个形状
    """
    if shape.has_text_frame:
        process_text_frame(shape.text_frame, stats)
    
    if shape.has_table:
        process_table(shape.table, stats)
    
    if hasattr(shape, 'shapes'):
        for sub_shape in shape.shapes:
            process_shape(sub_shape, stats)


def backup_file(file_path):
    """备份原始文件"""
    backup_path = f"{file_path}.backup"
    try:
        shutil.copy2(file_path, backup_path)
        print(f"ℹ️ 已备份原文件: {Path(backup_path).name}")
        return backup_path
    except Exception as e:
        print(f"⚠️ 备份文件失败: {e}")
        return None


def process_pptx(input_file):
    """
    处理PPTX文件
    """
    input_path = Path(input_file)
    
    if not input_path.exists():
        print(f"❌ 错误：文件不存在 - {input_file}")
        return False
    
    if input_path.suffix.lower() != '.pptx':
        print(f"❌ 错误：文件必须是.pptx格式")
        return False
    
    try:
        print(f"📖 正在读取文件: {input_path.name}")
        
        # 备份原文件
        backup_file(input_path)
        
        # 读取文档
        prs = Presentation(input_path)
        
        # 统计信息
        stats = {
            'quotes': 0,
            'punctuation': 0,
            'units': 0
        }
        
        total_slides = len(prs.slides)
        print(f"ℹ️ 文档包含 {total_slides} 张幻灯片")
        
        # 处理幻灯片母版
        print(f"🔄 正在处理幻灯片母版...")
        for slide_master in prs.slide_masters:
            for shape in slide_master.shapes:
                process_shape(shape, stats)
            for layout in slide_master.slide_layouts:
                for shape in layout.shapes:
                    process_shape(shape, stats)
        
        # 处理所有幻灯片
        print(f"🔄 正在处理幻灯片...")
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                process_shape(shape, stats)
        
        # 保存文件（覆盖原文件）
        print(f"💾 正在保存文件...")
        prs.save(input_path)
        
        print(f"✅ 处理完成！")
        print(f"   - 共替换了 {stats['quotes']} 个引号")
        print(f"   - 共替换了 {stats['punctuation']} 个标点符号")
        print(f"   - 共转换了 {stats['units']} 个单位")
        print(f"   - 输出文件: {input_path.name}")
        
        return True
        
    except Exception as e:
        print(f"❌ 处理失败：{e}")
        import traceback
        traceback.print_exc()
        return False


if __name__ == "__main__":
    # 获取输入文件（优先命令行参数，否则从 Finder 获取）
    files = get_input_files(sys.argv[1:], expected_ext='pptx', allow_multiple=False)
    
    if not files:
        print("❌ 错误：缺少文件名参数")
        print("\n使用方法：")
        print("    python3 text_formatter.py 文件名.pptx")
        print("    或在 Finder 中选择 .pptx 文件后运行")
        print("\n示例：")
        print("    python3 text_formatter.py presentation.pptx")
        sys.exit(1)
    
    input_file = files[0]
    
    print("=" * 50)
    print("文本格式自动修复工具 - PPTX版本")
    print("=" * 50)
    
    success = process_pptx(input_file)
    
    if success:
        print("\n🎉 全部完成！")
    else:
        print("\n❌ 处理失败，请检查错误信息")
        sys.exit(1)

