#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT表格样式工具

功能：
- 设置所有表格的 Header Row（标题行）
- 设置所有表格的 Banded Rows（镶边行）
- 设置所有表格的 First Column（首列）
- 自动备份原文件

用法:
    python3 table_style.py <input.pptx>
"""

import sys
import os
import shutil
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent.parent / "lib"))
from finder import get_input_files

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("❌ 错误: 缺少 python-pptx 库")
    print("💡 请运行: pip install python-pptx")
    sys.exit(1)


def show_message(msg_type, message):
    """显示格式化消息"""
    icons = {
        'success': '✅',
        'error': '❌',
        'warning': '⚠️',
        'info': 'ℹ️',
        'processing': '🔄'
    }
    icon = icons.get(msg_type, 'ℹ️')
    print(f"{icon} {message}")


def backup_file(file_path):
    """备份原始文件"""
    backup_path = f"{file_path}.backup"
    try:
        shutil.copy2(file_path, backup_path)
        show_message('info', f"已备份原文件: {os.path.basename(backup_path)}")
        return backup_path
    except Exception as e:
        show_message('warning', f"备份文件失败: {e}")
        return None


def set_table_style(table):
    """
    设置表格样式选项
    
    Args:
        table: pptx table对象
    
    Returns:
        bool: 是否成功设置
    """
    try:
        # Header Row - 标题行
        table.first_row = True
        
        # Banded Rows - 镶边行（交替行颜色）
        table.horz_banding = True
        
        # First Column - 首列
        table.first_col = True
        
        # 其他可选设置（默认关闭）
        # table.last_row = False      # Total Row - 汇总行
        # table.last_col = False      # Last Column - 末列
        # table.vert_banding = False  # Banded Columns - 镶边列
        
        return True
    except Exception as e:
        show_message('warning', f"设置表格样式失败: {e}")
        return False


def process_shape(shape, stats):
    """
    处理形状，查找表格
    
    Args:
        shape: pptx shape对象
        stats: 统计字典
    """
    # 处理表格
    if shape.has_table:
        if set_table_style(shape.table):
            stats['processed_tables'] += 1
    
    # 处理组合形状中的子形状
    if hasattr(shape, 'shapes'):
        for sub_shape in shape.shapes:
            process_shape(sub_shape, stats)


def process_presentation(input_path, output_path=None):
    """
    处理PPT文档中所有表格的样式
    
    Args:
        input_path: 输入文件路径
        output_path: 输出文件路径，如果为None则覆盖原文件
    
    Returns:
        bool: 是否成功
    """
    try:
        # 验证输入文件
        if not os.path.exists(input_path):
            show_message('error', f"文件不存在: {input_path}")
            return False
        
        if not input_path.lower().endswith('.pptx'):
            show_message('error', "只支持.pptx格式的文件")
            return False
        
        show_message('processing', f"正在处理文件: {os.path.basename(input_path)}")
        
        # 备份原文件
        backup_file(input_path)
        
        # 打开PPT
        prs = Presentation(input_path)
        
        # 统计信息
        total_slides = len(prs.slides)
        stats = {'processed_tables': 0}
        
        show_message('info', f"文档包含 {total_slides} 张幻灯片")
        
        # 处理所有幻灯片
        show_message('processing', "正在处理表格样式...")
        for i, slide in enumerate(prs.slides, 1):
            try:
                for shape in slide.shapes:
                    process_shape(shape, stats)
            except Exception as e:
                show_message('warning', f"处理第{i}张幻灯片时出错: {e}")
                continue
        
        if stats['processed_tables'] > 0:
            show_message('info', f"已处理 {stats['processed_tables']} 个表格")
        else:
            show_message('warning', "未找到任何表格")
        
        # 保存文档
        output_file = output_path if output_path else input_path
        prs.save(output_file)
        
        show_message('success', f"表格样式设置完成: {os.path.basename(output_file)}")
        show_message('info', "已启用: Header Row, Banded Rows, First Column")
        
        return True
        
    except Exception as e:
        show_message('error', f"处理文件时出错: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """主函数"""
    # 获取输入文件（优先命令行参数，否则从 Finder 获取）
    files = get_input_files(sys.argv[1:], expected_ext='pptx', allow_multiple=False)
    
    if not files:
        print("PPT表格样式工具")
        print()
        print("用法:")
        print(f"  python {os.path.basename(__file__)} <输入文件.pptx>")
        print("  或在 Finder 中选择 .pptx 文件后运行")
        print()
        print("功能:")
        print("  为所有表格启用以下样式选项:")
        print("  ✓ Header Row（标题行）")
        print("  ✓ Banded Rows（镶边行）")
        print("  ✓ First Column（首列）")
        print()
        print("示例:")
        print(f"  python {os.path.basename(__file__)} presentation.pptx")
        sys.exit(1)
    
    input_path = files[0]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None
    
    success = process_presentation(input_path, output_path)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()

