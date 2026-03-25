#!/usr/bin/env python3
"""
PPTX 文档标准化工具集 v3.0.0

将 4 个独立 PPTX 脚本合并为一个多子命令工具，直接函数调用，不再依赖 subprocess。

子命令:
    font    - 字体统一为微软雅黑（处理所有文本框、表格、母版）
    format  - 文本格式修复（引号、标点、中文单位转标准符号）
    table   - 表格样式设置（标题行、镶边行、首列）
    all     - 一键标准化：依次执行 format -> font -> table

用法:
    python3 pptx_tools.py <subcommand> [file...]
    python3 pptx_tools.py font presentation.pptx
    python3 pptx_tools.py format file1.pptx file2.pptx
    python3 pptx_tools.py table presentation.pptx
    python3 pptx_tools.py all presentation.pptx

作者: tianli
版本: 3.0.0
日期: 2026-03-25
"""

import argparse
import os
import re
import shutil
import sys
import traceback
from pathlib import Path

# ── lib 路径 ──────────────────────────────────────────────────────────
sys.path.insert(0, str(Path(__file__).parent.parent.parent / "lib"))
from finder import get_input_files
from progress import ProgressTracker
from text_format_rules import fix_punctuation, fix_quotes, fix_units

# ── 第三方依赖 ────────────────────────────────────────────────────────
try:
    from lxml import etree
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from pptx.util import Pt
except ImportError:
    print("❌ 错误: 缺少 python-pptx 或 lxml 库")
    print("💡 请运行: pip install python-pptx lxml")
    sys.exit(1)


# =====================================================================
#  共用工具函数
# =====================================================================

def show_message(msg_type, message):
    """显示格式化消息"""
    icons = {"success": "✅", "error": "❌", "warning": "⚠️", "info": "ℹ️", "processing": "🔄"}
    icon = icons.get(msg_type, "ℹ️")
    print(f"{icon} {message}")


def backup_file(file_path):
    """备份原始文件"""
    backup_path = f"{file_path}.backup"
    try:
        shutil.copy2(file_path, backup_path)
        show_message("info", f"已备份原文件: {Path(backup_path).name}")
        return backup_path
    except Exception as e:
        show_message("warning", f"备份文件失败: {e}")
        return None


# =====================================================================
#  子命令: font — 字体统一为微软雅黑
# =====================================================================

# 目标字体
TARGET_FONT = "Microsoft YaHei"

# XML 命名空间
NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def font_set_for_run(run, font_name=TARGET_FONT):
    """
    为 run 设置字体（强制 XML 级别设置）

    Args:
        run: pptx run 对象
        font_name: 字体名称
    """
    try:
        # 1. 使用 API 设置
        run.font.name = font_name

        # 2. 强制 XML 级别设置 - 直接操作 rPr 元素
        rPr = run._r.get_or_add_rPr()

        # 查找或创建 a:latin 元素（西文字体）
        latin = rPr.find(qn("a:latin"))
        if latin is None:
            latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", font_name)

        # 查找或创建 a:ea 元素（东亚字体 - 中文）
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font_name)

        # 查找或创建 a:cs 元素（复杂脚本字体）
        cs = rPr.find(qn("a:cs"))
        if cs is None:
            cs = etree.SubElement(rPr, qn("a:cs"))
        cs.set("typeface", font_name)

    except Exception:
        # 某些 run 可能没有字体属性
        pass


def font_set_paragraph_default(paragraph, font_name=TARGET_FONT):
    """
    设置段落的默认字体属性（defRPr）

    Args:
        paragraph: pptx paragraph 对象
        font_name: 字体名称
    """
    try:
        pPr = paragraph._p.get_or_add_pPr()

        # 查找或创建 defRPr（默认文本属性）
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn("a:defRPr"))

        # 设置 latin 字体
        latin = defRPr.find(qn("a:latin"))
        if latin is None:
            latin = etree.SubElement(defRPr, qn("a:latin"))
        latin.set("typeface", font_name)

        # 设置 ea 字体（东亚）
        ea = defRPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(defRPr, qn("a:ea"))
        ea.set("typeface", font_name)

        # 设置 cs 字体
        cs = defRPr.find(qn("a:cs"))
        if cs is None:
            cs = etree.SubElement(defRPr, qn("a:cs"))
        cs.set("typeface", font_name)

    except Exception:
        pass


def font_set_endParaRPr(paragraph, font_name=TARGET_FONT):
    """
    设置段落结束符的字体属性（endParaRPr）

    Args:
        paragraph: pptx paragraph 对象
        font_name: 字体名称
    """
    try:
        endParaRPr = paragraph._p.find(qn("a:endParaRPr"))
        if endParaRPr is not None:
            # 设置 latin 字体
            latin = endParaRPr.find(qn("a:latin"))
            if latin is None:
                latin = etree.SubElement(endParaRPr, qn("a:latin"))
            latin.set("typeface", font_name)

            # 设置 ea 字体
            ea = endParaRPr.find(qn("a:ea"))
            if ea is None:
                ea = etree.SubElement(endParaRPr, qn("a:ea"))
            ea.set("typeface", font_name)

            # 设置 cs 字体
            cs = endParaRPr.find(qn("a:cs"))
            if cs is None:
                cs = etree.SubElement(endParaRPr, qn("a:cs"))
            cs.set("typeface", font_name)
    except Exception:
        pass


def font_process_text_frame(text_frame, stats):
    """
    处理文本框中的所有段落和 run（字体设置）

    Args:
        text_frame: pptx text_frame 对象
        stats: 统计字典
    """
    for paragraph in text_frame.paragraphs:
        # 设置段落默认字体
        font_set_paragraph_default(paragraph, TARGET_FONT)
        # 设置段落结束符字体
        font_set_endParaRPr(paragraph, TARGET_FONT)

        # 处理每个 run
        for run in paragraph.runs:
            font_set_for_run(run, TARGET_FONT)
            stats["font_processed_runs"] += 1


def font_process_table(table, stats):
    """
    处理表格中的所有单元格（字体设置）

    Args:
        table: pptx table 对象
        stats: 统计字典
    """
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                font_process_text_frame(cell.text_frame, stats)
                stats["font_processed_tables"] += 1


def font_process_shape(shape, stats):
    """
    处理单个形状（字体设置）

    Args:
        shape: pptx shape 对象
        stats: 统计字典
    """
    # 处理有文本框的形状
    if shape.has_text_frame:
        font_process_text_frame(shape.text_frame, stats)
        stats["font_processed_shapes"] += 1

    # 处理表格
    if shape.has_table:
        font_process_table(shape.table, stats)

    # 处理组合形状中的子形状
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            font_process_shape(sub_shape, stats)


def font_process_slide(slide, stats):
    """
    处理单个幻灯片（字体设置）

    Args:
        slide: pptx slide 对象
        stats: 统计字典
    """
    for shape in slide.shapes:
        font_process_shape(shape, stats)


def font_process_slide_master(slide_master, stats):
    """
    处理幻灯片母版（字体设置）

    Args:
        slide_master: pptx slide_master 对象
        stats: 统计字典
    """
    # 处理母版中的形状
    for shape in slide_master.shapes:
        font_process_shape(shape, stats)

    # 处理母版的布局
    for layout in slide_master.slide_layouts:
        for shape in layout.shapes:
            font_process_shape(shape, stats)


def font_process_presentation(input_path, do_backup=True):
    """
    格式化 PPT 文档中所有文字的字体为微软雅黑

    Args:
        input_path: 输入文件路径
        do_backup: 是否备份原文件

    Returns:
        bool: 是否成功
    """
    try:
        # 验证输入文件
        if not os.path.exists(input_path):
            show_message("error", f"文件不存在: {input_path}")
            return False

        if not input_path.lower().endswith(".pptx"):
            show_message("error", "只支持.pptx格式的文件")
            return False

        show_message("processing", f"正在处理文件: {os.path.basename(input_path)}")

        # 备份原文件
        if do_backup:
            backup_file(input_path)

        # 打开 PPT
        prs = Presentation(input_path)

        # 统计信息
        total_slides = len(prs.slides)
        stats = {
            "font_processed_shapes": 0,
            "font_processed_runs": 0,
            "font_processed_tables": 0,
        }

        show_message("info", f"文档包含 {total_slides} 张幻灯片")

        # 处理幻灯片母版（重要：这里的字体设置会影响整个 PPT）
        show_message("processing", "正在处理幻灯片母版...")
        for slide_master in prs.slide_masters:
            try:
                font_process_slide_master(slide_master, stats)
            except Exception as e:
                show_message("warning", f"处理母版时出错: {e}")

        # 处理所有幻灯片
        show_message("processing", "正在处理幻灯片...")
        for i, slide in enumerate(prs.slides, 1):
            try:
                font_process_slide(slide, stats)
            except Exception as e:
                show_message("warning", f"处理第{i}张幻灯片时出错: {e}")
                continue

        show_message(
            "info",
            f"已处理 {stats['font_processed_shapes']} 个形状, "
            f"{stats['font_processed_runs']} 个文本run",
        )
        if stats["font_processed_tables"] > 0:
            show_message("info", f"已处理 {stats['font_processed_tables']} 个表格单元格")

        # 保存文档
        prs.save(input_path)

        show_message("success", f"字体格式化完成: {os.path.basename(input_path)}")
        show_message("info", f"所有文字已设置为: {TARGET_FONT}")

        return True

    except Exception as e:
        show_message("error", f"处理文件时出错: {e}")
        traceback.print_exc()
        return False


# =====================================================================
#  子命令: format — 文本格式修复（引号、标点、单位）
# =====================================================================

def format_process_text(text, stats):
    """
    处理文本，应用所有文本转换
    """
    if not text:
        return text

    result, quote_count = fix_quotes(text)
    result, punct_count = format_fix_punctuation(result)
    result, unit_count = format_fix_units(result)

    stats["format_quotes"] += quote_count
    stats["format_punctuation"] += punct_count
    stats["format_units"] += unit_count

    return result


def format_process_run(run, stats):
    """
    处理单个 run 的文本（格式修复）
    """
    if run.text:
        original = run.text
        fixed = format_process_text(original, stats)
        if fixed != original:
            run.text = fixed


def format_process_text_frame(text_frame, stats):
    """
    处理文本框中的所有段落和 run（格式修复）
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            format_process_run(run, stats)


def format_process_table(table, stats):
    """
    处理表格中的所有单元格（格式修复）
    """
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                format_process_text_frame(cell.text_frame, stats)


def format_process_shape(shape, stats):
    """
    处理单个形状（格式修复）
    """
    if shape.has_text_frame:
        format_process_text_frame(shape.text_frame, stats)

    if shape.has_table:
        format_process_table(shape.table, stats)

    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            format_process_shape(sub_shape, stats)


def format_process_presentation(input_path, do_backup=True):
    """
    处理 PPTX 文件的文本格式（引号、标点、单位）

    Args:
        input_path: 输入文件路径
        do_backup: 是否备份原文件

    Returns:
        bool: 是否成功
    """
    input_p = Path(input_path)

    if not input_p.exists():
        show_message("error", f"文件不存在: {input_path}")
        return False

    if input_p.suffix.lower() != ".pptx":
        show_message("error", "文件必须是.pptx格式")
        return False

    try:
        show_message("processing", f"正在处理文件: {input_p.name}")

        # 备份原文件
        if do_backup:
            backup_file(input_path)

        # 读取文档
        prs = Presentation(input_path)

        # 统计信息
        stats = {"format_quotes": 0, "format_punctuation": 0, "format_units": 0}

        total_slides = len(prs.slides)
        show_message("info", f"文档包含 {total_slides} 张幻灯片")

        # 处理幻灯片母版
        show_message("processing", "正在处理幻灯片母版...")
        for slide_master in prs.slide_masters:
            for shape in slide_master.shapes:
                format_process_shape(shape, stats)
            for layout in slide_master.slide_layouts:
                for shape in layout.shapes:
                    format_process_shape(shape, stats)

        # 处理所有幻灯片
        show_message("processing", "正在处理幻灯片...")
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                format_process_shape(shape, stats)

        # 保存文件（覆盖原文件）
        prs.save(input_path)

        show_message("success", "文本格式修复完成！")
        show_message("info", f"   共替换了 {stats['format_quotes']} 个引号")
        show_message("info", f"   共替换了 {stats['format_punctuation']} 个标点符号")
        show_message("info", f"   共转换了 {stats['format_units']} 个单位")

        return True

    except Exception as e:
        show_message("error", f"处理失败: {e}")
        traceback.print_exc()
        return False


# =====================================================================
#  子命令: table — 表格样式设置
# =====================================================================

def table_set_style(table):
    """
    设置表格样式选项

    Args:
        table: pptx table 对象

    Returns:
        bool: 是否成功设置
    """
    try:
        # Header Row - 标题行
        table.first_row = True

        # Banded Rows - 镶边行（交替行颜色）
        table.horz_banding = True

        # First Column - 首列
        table.first_col = True

        # 其他可选设置（默认关闭）
        # table.last_row = False      # Total Row - 汇总行
        # table.last_col = False      # Last Column - 末列
        # table.vert_banding = False  # Banded Columns - 镶边列

        return True
    except Exception as e:
        show_message("warning", f"设置表格样式失败: {e}")
        return False


def table_process_shape(shape, stats):
    """
    处理形状，查找表格（表格样式设置）

    Args:
        shape: pptx shape 对象
        stats: 统计字典
    """
    # 处理表格
    if shape.has_table:
        if table_set_style(shape.table):
            stats["table_processed"] += 1

    # 处理组合形状中的子形状
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            table_process_shape(sub_shape, stats)


def table_process_presentation(input_path, do_backup=True):
    """
    处理 PPT 文档中所有表格的样式

    Args:
        input_path: 输入文件路径
        do_backup: 是否备份原文件

    Returns:
        bool: 是否成功
    """
    try:
        # 验证输入文件
        if not os.path.exists(input_path):
            show_message("error", f"文件不存在: {input_path}")
            return False

        if not input_path.lower().endswith(".pptx"):
            show_message("error", "只支持.pptx格式的文件")
            return False

        show_message("processing", f"正在处理文件: {os.path.basename(input_path)}")

        # 备份原文件
        if do_backup:
            backup_file(input_path)

        # 打开 PPT
        prs = Presentation(input_path)

        # 统计信息
        total_slides = len(prs.slides)
        stats = {"table_processed": 0}

        show_message("info", f"文档包含 {total_slides} 张幻灯片")

        # 处理所有幻灯片
        show_message("processing", "正在处理表格样式...")
        for i, slide in enumerate(prs.slides, 1):
            try:
                for shape in slide.shapes:
                    table_process_shape(shape, stats)
            except Exception as e:
                show_message("warning", f"处理第{i}张幻灯片时出错: {e}")
                continue

        if stats["table_processed"] > 0:
            show_message("info", f"已处理 {stats['table_processed']} 个表格")
        else:
            show_message("warning", "未找到任何表格")

        # 保存文档
        prs.save(input_path)

        show_message("success", f"表格样式设置完成: {os.path.basename(input_path)}")
        show_message("info", "已启用: Header Row, Banded Rows, First Column")

        return True

    except Exception as e:
        show_message("error", f"处理文件时出错: {e}")
        traceback.print_exc()
        return False


# =====================================================================
#  子命令: all — 一键标准化（format -> font -> table）
# =====================================================================

def all_process_presentation(input_path):
    """
    应用所有 PPTX 标准化处理（直接函数调用，不再使用 subprocess）

    执行顺序：
    1. 文本格式修复（引号、标点、单位）
    2. 字体统一为微软雅黑
    3. 表格样式设置

    Args:
        input_path: 输入文件路径

    Returns:
        bool: 是否全部成功
    """
    input_p = Path(input_path)

    # 检查文件是否存在
    if not input_p.exists():
        show_message("error", f"文件不存在: {input_path}")
        return False

    if input_p.suffix.lower() != ".pptx":
        show_message("error", "只支持 .pptx 文件")
        return False

    print("=" * 70)
    print("🚀 开始 PPT 文档标准化处理")
    print("=" * 70)
    print(f"📄 文件: {input_p.name}")
    print()

    # 先备份一次（后续步骤不再重复备份）
    backup_file(input_path)

    # 执行顺序：format -> font -> table
    steps = [
        {
            "name": "步骤 1/3: 文本格式修复",
            "func": format_process_presentation,
        },
        {
            "name": "步骤 2/3: 字体统一为微软雅黑",
            "func": font_process_presentation,
        },
        {
            "name": "步骤 3/3: 表格样式设置",
            "func": table_process_presentation,
        },
    ]

    success_count = 0
    failed_steps = []

    for step in steps:
        print("\n" + "=" * 70)
        print(f"▶️  {step['name']}")
        print("=" * 70)

        # 调用处理函数，do_backup=False 因为已经在上面统一备份过了
        if step["func"](str(input_p), do_backup=False):
            success_count += 1
            print(f"✅ {step['name']} 完成")
        else:
            failed_steps.append(step["name"])
            print(f"⚠️ {step['name']} 失败（继续执行后续步骤）")

    # 总结
    print("\n" + "=" * 70)
    print("📊 处理总结")
    print("=" * 70)
    print(f"✅ 成功: {success_count}/{len(steps)} 个步骤")

    if failed_steps:
        print(f"⚠️ 失败: {len(failed_steps)} 个步骤")
        for step_name in failed_steps:
            print(f"   - {step_name}")
    else:
        print("🎉 所有步骤执行成功！")

    print(f"\n📄 处理完成: {input_p.name}")
    print("=" * 70)

    return len(failed_steps) == 0


# =====================================================================
#  CLI 入口
# =====================================================================

def build_parser():
    """构建 argparse 解析器"""
    parser = argparse.ArgumentParser(
        prog="pptx_tools",
        description="PPTX 文档标准化工具集",
        epilog=(
            "子命令说明:\n"
            "  font    字体统一为微软雅黑\n"
            "  format  文本格式修复（引号、标点、单位）\n"
            "  table   表格样式（标题行、镶边行、首列）\n"
            "  all     一键标准化: format -> font -> table\n"
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "subcommand",
        choices=["font", "format", "table", "all"],
        help="子命令: font | format | table | all",
    )
    parser.add_argument(
        "files",
        nargs="*",
        help="PPTX 文件路径（可多个；不提供则从 Finder 选中获取）",
    )
    return parser


def main():
    parser = build_parser()

    # 没有参数时显示帮助
    if len(sys.argv) < 2:
        parser.print_help()
        sys.exit(1)

    args = parser.parse_args()

    # 获取输入文件
    files = get_input_files(args.files, expected_ext="pptx")

    if not files:
        show_message("error", "未找到 .pptx 文件")
        print("\n用法: python3 pptx_tools.py <subcommand> [file...]")
        print("  或在 Finder 中选择 .pptx 文件后运行")
        sys.exit(1)

    # 分发子命令
    dispatch = {
        "font": font_process_presentation,
        "format": format_process_presentation,
        "table": table_process_presentation,
        "all": all_process_presentation,
    }

    handler = dispatch[args.subcommand]
    tracker = ProgressTracker()

    for file_path in files:
        print(f"\n{'=' * 50}")
        print(f"处理文件: {Path(file_path).name}")
        print("=" * 50)

        file_str = str(file_path)

        if args.subcommand == "all":
            success = handler(file_str)
        else:
            success = handler(file_str, do_backup=True)

        if success:
            tracker.add_success()
        else:
            tracker.add_error()

    print(f"\n{'=' * 50}")
    tracker.show_summary("文件处理")


if __name__ == "__main__":
    main()
