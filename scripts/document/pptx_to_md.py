#!/usr/bin/env python3
"""
PPTX转Markdown转换工具 - 将PPTX演示文稿转换为Markdown格式
版本: 2.0.0
作者: tianli
"""

import sys
import argparse
from pathlib import Path
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).parent.parent.parent / "lib"))
from display import show_success, show_error, show_warning, show_info, show_processing
from file_ops import (validate_input_file, ensure_directory, fatal_error,
                      check_python_packages, find_files_by_extension, get_file_basename)
from finder import get_input_files
from progress import ProgressTracker

SCRIPT_VERSION = "2.0.0"

def check_dependencies():
    show_info("检查依赖项...")
    if not check_python_packages('pptx'):
        sys.exit(1)
    show_success("依赖检查完成")

def convert_pptx_to_md_single(file_path: Path, output_dir: Path) -> bool:
    if not validate_input_file(file_path):
        return False

    if file_path.suffix.lower() != '.pptx':
        show_warning(f"跳过非PPTX文件: {file_path.name}")
        return False

    base_name = get_file_basename(file_path)
    # 直接在原文件目录生成 .md 文件
    output_file = file_path.parent / f"{base_name}.md"
    
    show_processing(f"转换 {file_path.name} 为 Markdown...")

    try:
        prs = Presentation(file_path)
        with open(output_file, 'w', encoding='utf-8') as md_file:
            for i, slide in enumerate(prs.slides, 1):
                md_file.write(f"## Slide {i}\n\n")
                
                notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        md_file.write(shape.text + '\n\n')

                if notes:
                    md_file.write(f"### Speaker Notes\n\n{notes}\n\n")
                
                md_file.write("---\n\n")
        
        show_success(f"成功转换: {file_path.name} -> {output_file.name}")
        return True
    except Exception as e:
        show_error(f"转换失败 {file_path.name}: {e}")
        return False

def collect_pptx_files(input_paths: list, recursive: bool = False) -> list:
    """从输入路径列表中收集所有PPTX文件"""
    all_files = []
    
    for input_path in input_paths:
        path_obj = Path(input_path)
        
        if path_obj.is_file():
            # 如果是文件，直接检查扩展名
            if path_obj.suffix.lower() == '.pptx':
                all_files.append(path_obj)
            else:
                show_warning(f"跳过非PPTX文件: {path_obj.name}")
        elif path_obj.is_dir():
            # 如果是目录，查找其中的PPTX文件
            found_files = find_files_by_extension(path_obj, 'pptx', recursive)
            all_files.extend(found_files)
        else:
            show_error(f"路径不存在: {input_path}")
    
    return all_files

def main():
    # 无参数时从 Finder 获取选中的文件
    if len(sys.argv) == 1:
        files = get_input_files([], expected_ext='pptx')
        if files:
            sys.argv.extend(files)
    
    parser = argparse.ArgumentParser(description="PPTX转Markdown转换工具")
    parser.add_argument("input_paths", nargs='+', help="一个或多个PPTX文件/目录路径")
    parser.add_argument("-o", "--output", help="输出目录 (默认: ./converted_md)")
    parser.add_argument("-r", "--recursive", action="store_true", help="递归处理目录")
    parser.add_argument('--version', action='version', version=f'%(prog)s {SCRIPT_VERSION}')
    args = parser.parse_args()

    check_dependencies()

    files_to_process = collect_pptx_files(args.input_paths, args.recursive)

    if not files_to_process:
        show_warning("未找到任何PPTX文件")
        sys.exit(0)

    total_success = 0
    progress = ProgressTracker()

    for file_path in files_to_process:
        show_processing(f"处理 {file_path.name}")
        if convert_pptx_to_md_single(file_path, None):
            total_success += 1
            progress.add_success()
        else:
            progress.add_failure()
            
    show_info("\n处理完成")
    show_success(f"总共成功转换了 {total_success} 个文件")

if __name__ == "__main__":
    main()

