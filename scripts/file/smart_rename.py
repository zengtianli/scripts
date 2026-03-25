#!/usr/bin/env python3
"""Downloads 智能重命名工具 - AI 驱动的文件分析、分组与重命名。

用法:
    python3 smart_rename.py analyze --all             # 分析所有分类目录
    python3 smart_rename.py analyze --dir ~/Downloads/文档  # 分析指定目录
    python3 smart_rename.py analyze --all --dry-run    # 只分析前 5 个文件（测试）
    python3 smart_rename.py execute                    # 执行审核后的重命名计划
    python3 smart_rename.py rollback                   # 回滚上次执行
"""

import argparse
import difflib
import hashlib
import json
import logging
import re
import shutil
import sys
from collections import defaultdict
from datetime import datetime
from pathlib import Path

import yaml

# Add scripts root to path for importing shared modules
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from tools.llm_client import chat as llm_chat

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
SCRIPT_DIR = Path(__file__).parent
CONFIG_FILE = SCRIPT_DIR / "smart_rename_config.yaml"

SYSTEM_PROMPT = """\
你是文件命名助手。分析以下文件列表，完成三件事：
1. 按项目/主题分组（group 字段用简短中文，2-6字）
2. 识别同一文件的不同版本（版本号较旧的 action 设为 archive）
3. 为每个文件建议新名称，格式：项目名-内容描述-日期.扩展名

规则：
- 日期从文件名或内容中提取，格式 YYYYMMDD。找不到日期则省略日期部分
- 项目名用简短中文（2-6字），内容描述用简短中文（2-8字）
- 如果当前文件名已经足够清晰，new_name 可以只做微调（去掉前缀0、规范化日期格式等）
- 完全重复的文件（我会标注 MD5 相同）action 设为 delete
- 不确定的文件 action 设为 skip

只返回 JSON 数组，不要其他文字：
[{"index": 1, "group": "项目名", "new_name": "建议文件名.扩展名", "action": "rename", "reason": "简短说明"}]

action 可选值：rename（重命名）、archive（归档旧版本）、delete（删除重复）、skip（跳过）"""

# Document extensions that support content extraction
CONTENT_EXTRACTABLE = {
    ".docx", ".doc", ".pdf", ".xlsx", ".xls",
    ".pptx", ".ppt", ".txt", ".md", ".html", ".csv",
}


# ---------------------------------------------------------------------------
# Config & Logging
# ---------------------------------------------------------------------------
def load_config() -> dict:
    with open(CONFIG_FILE, encoding="utf-8") as f:
        return yaml.safe_load(f)


def setup_logging(log_file: str) -> logging.Logger:
    log_path = Path(log_file).expanduser()
    log_path.parent.mkdir(parents=True, exist_ok=True)
    logger = logging.getLogger("smart_rename")
    logger.setLevel(logging.INFO)
    fmt = logging.Formatter("%(asctime)s  %(message)s", datefmt="%Y-%m-%d %H:%M:%S")
    fh = logging.FileHandler(log_path, encoding="utf-8")
    fh.setFormatter(fmt)
    logger.addHandler(fh)
    sh = logging.StreamHandler(sys.stdout)
    sh.setFormatter(fmt)
    logger.addHandler(sh)
    return logger


# ---------------------------------------------------------------------------
# File Scanning
# ---------------------------------------------------------------------------
def compute_md5(filepath: Path, chunk_size: int = 8192) -> str:
    """Compute MD5 of first 8KB for quick comparison."""
    h = hashlib.md5()
    with open(filepath, "rb") as f:
        data = f.read(chunk_size)
        h.update(data)
    return h.hexdigest()


def extract_content_preview(filepath: Path, max_chars: int = 500) -> str:
    """Extract text preview from document files."""
    ext = filepath.suffix.lower()
    try:
        if ext == ".docx":
            return _extract_docx(filepath, max_chars)
        elif ext == ".pdf":
            return _extract_pdf(filepath, max_chars)
        elif ext == ".xlsx":
            return _extract_xlsx(filepath, max_chars)
        elif ext == ".pptx":
            return _extract_pptx(filepath, max_chars)
        elif ext in (".txt", ".md", ".html", ".csv"):
            return _extract_text(filepath, max_chars)
    except Exception as e:
        return f"[提取失败: {e}]"
    return ""


def _extract_docx(filepath: Path, max_chars: int) -> str:
    from docx import Document
    doc = Document(str(filepath))
    texts = []
    for para in doc.paragraphs[:5]:
        if para.text.strip():
            texts.append(para.text.strip())
        if sum(len(t) for t in texts) >= max_chars:
            break
    return "\n".join(texts)[:max_chars]


def _extract_pdf(filepath: Path, max_chars: int) -> str:
    import pdfplumber
    with pdfplumber.open(str(filepath)) as pdf:
        if not pdf.pages:
            return ""
        text = pdf.pages[0].extract_text() or ""
        return text[:max_chars]


def _extract_xlsx(filepath: Path, max_chars: int) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(filepath), read_only=True, data_only=True)
    parts = [f"工作表: {', '.join(wb.sheetnames)}"]
    ws = wb.active
    if ws:
        for row in ws.iter_rows(max_row=3, values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)[:max_chars]


def _extract_pptx(filepath: Path, max_chars: int) -> str:
    from pptx import Presentation
    prs = Presentation(str(filepath))
    titles = []
    for slide in prs.slides[:5]:
        if slide.shapes.title and slide.shapes.title.text.strip():
            titles.append(slide.shapes.title.text.strip())
    return "\n".join(titles)[:max_chars]


def _extract_text(filepath: Path, max_chars: int) -> str:
    encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
    for enc in encodings:
        try:
            with open(filepath, encoding=enc) as f:
                return f.read(max_chars)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return ""


def scan_files(
    directories: list[Path],
    ignore_prefixes: list[str],
    max_chars: int,
    logger: logging.Logger,
    dry_run_limit: int = 0,
) -> list[dict]:
    """Scan directories and collect file metadata."""
    files = []
    for d in directories:
        if not d.exists():
            logger.warning("目录不存在，跳过: %s", d)
            continue
        logger.info("扫描: %s", d)
        for item in sorted(d.iterdir()):
            if not item.is_file():
                continue
            if any(item.name.startswith(p) for p in ignore_prefixes):
                continue
            ext = item.suffix.lower()
            can_extract = ext in CONTENT_EXTRACTABLE
            preview = ""
            if can_extract:
                logger.info("  提取内容: %s", item.name)
                preview = extract_content_preview(item, max_chars)

            files.append({
                "path": str(item),
                "dir": item.parent.name,
                "name": item.name,
                "ext": ext,
                "size": item.stat().st_size,
                "mtime": datetime.fromtimestamp(item.stat().st_mtime).strftime("%Y-%m-%d"),
                "md5": compute_md5(item),
                "content_preview": preview,
            })
            if dry_run_limit and len(files) >= dry_run_limit:
                return files
    return files


# ---------------------------------------------------------------------------
# Grouping
# ---------------------------------------------------------------------------
def find_duplicates(files: list[dict]) -> dict[str, list[int]]:
    """Find files with identical MD5. Returns md5 -> list of indices."""
    md5_map = defaultdict(list)
    for i, f in enumerate(files):
        md5_map[f["md5"]].append(i)
    return {k: v for k, v in md5_map.items() if len(v) > 1}


def find_similar_names(files: list[dict], threshold: float = 0.7) -> list[list[int]]:
    """Find groups of files with similar names using SequenceMatcher."""
    n = len(files)
    visited = set()
    groups = []

    for i in range(n):
        if i in visited:
            continue
        group = [i]
        for j in range(i + 1, n):
            if j in visited:
                continue
            if files[i]["ext"] != files[j]["ext"]:
                continue
            if files[i]["dir"] != files[j]["dir"]:
                continue
            ratio = difflib.SequenceMatcher(
                None, files[i]["name"], files[j]["name"]
            ).ratio()
            if ratio >= threshold:
                group.append(j)
                visited.add(j)
        if len(group) > 1:
            visited.add(i)
            groups.append(group)

    return groups


# ---------------------------------------------------------------------------
# AI Analysis
# ---------------------------------------------------------------------------
def build_batch_prompt(files: list[dict], indices: list[int], dup_indices: set[int]) -> str:
    """Build user message for a batch of files.

    Uses 1-based sequential numbering within the batch (not global indices)
    to avoid AI renumbering issues.
    """
    lines = ["文件列表："]
    for batch_pos, idx in enumerate(indices, 1):
        f = files[idx]
        size_kb = f["size"] / 1024
        line = f'{batch_pos}. 文件名: {f["name"]} | 大小: {size_kb:.0f}KB | 修改: {f["mtime"]} | 目录: {f["dir"]}'
        if idx in dup_indices:
            line += " | [MD5重复]"
        if f["content_preview"]:
            preview = f["content_preview"][:300].replace("\n", " ")
            line += f" | 内容预览: {preview}"
        lines.append(line)
    return "\n".join(lines)


def parse_ai_response(text: str) -> list[dict]:
    """Parse JSON array from AI response, tolerant of surrounding text."""
    # Try to find JSON array in the response
    match = re.search(r"\[.*\]", text, re.DOTALL)
    if match:
        try:
            return json.loads(match.group())
        except json.JSONDecodeError:
            pass
    return []


def analyze_with_ai(
    files: list[dict],
    config: dict,
    logger: logging.Logger,
) -> list[dict]:
    """Send files in batches to AI for analysis. Returns merged results."""
    ai_cfg = config["ai"]
    batch_size = ai_cfg["batch_size"]

    # Find duplicate indices for annotation
    dups = find_duplicates(files)
    dup_indices = set()
    for indices in dups.values():
        dup_indices.update(indices)

    all_indices = list(range(len(files)))
    results = [None] * len(files)

    # Process in batches
    for batch_start in range(0, len(all_indices), batch_size):
        batch = all_indices[batch_start: batch_start + batch_size]
        batch_num = batch_start // batch_size + 1
        total_batches = (len(all_indices) + batch_size - 1) // batch_size
        logger.info("AI 分析 batch %d/%d (%d 个文件)...", batch_num, total_batches, len(batch))

        prompt = build_batch_prompt(files, batch, dup_indices)
        try:
            response = llm_chat(system=SYSTEM_PROMPT, message=prompt)
            items = parse_ai_response(response)
            for item in items:
                batch_pos = item.get("index", 0) - 1  # 1-based to 0-based within batch
                if 0 <= batch_pos < len(batch):
                    global_idx = batch[batch_pos]
                    item["index"] = global_idx + 1  # Store global 1-based index
                    results[global_idx] = item
        except Exception as e:
            logger.error("  AI 调用失败: %s", e)
            # Fill failed batch with skip
            for idx in batch:
                if results[idx] is None:
                    results[idx] = {
                        "index": idx + 1,
                        "group": "未分析",
                        "new_name": files[idx]["name"],
                        "action": "skip",
                        "reason": f"AI 调用失败: {e}",
                    }

        # Brief pause between batches
        if batch_start + batch_size < len(all_indices):
            import time
            time.sleep(1)

    # Fill any None results
    for i, r in enumerate(results):
        if r is None:
            results[i] = {
                "index": i + 1,
                "group": "未分组",
                "new_name": files[i]["name"],
                "action": "skip",
                "reason": "未被AI分析",
            }

    return results


# ---------------------------------------------------------------------------
# Plan Generation
# ---------------------------------------------------------------------------
def generate_plan_md(
    files: list[dict],
    ai_results: list[dict],
    plan_path: Path,
    logger: logging.Logger,
) -> None:
    """Generate the rename plan as a Markdown file."""
    # Group by AI-assigned group, then by directory
    groups = defaultdict(list)
    for i, (f, r) in enumerate(zip(files, ai_results)):
        key = (f["dir"], r.get("group", "未分组"))
        groups[key].append((i, f, r))

    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    dir_counts = defaultdict(int)
    for f in files:
        dir_counts[f["dir"]] += 1

    lines = [
        "# Downloads 智能重命名建议\n",
        f"> 生成时间: {now}",
        f"> 扫描范围: {', '.join(f'{d}/({c})' for d, c in sorted(dir_counts.items()))}",
        "> 使用说明:",
        "> - 编辑此文件后运行 `python3 smart_rename.py execute`",
        "> - 修改「操作」列：rename / archive / delete / skip",
        "> - 修改「建议新名」列可自定义文件名",
        "> - 删除整行 = 跳过该文件",
        "",
    ]

    # Sort groups by directory, then group name
    sorted_keys = sorted(groups.keys())
    current_dir = None

    for dir_name, group_name in sorted_keys:
        if dir_name != current_dir:
            current_dir = dir_name
            lines.append(f"## {dir_name}/ ({dir_counts[dir_name]} 个文件)\n")

        items = groups[(dir_name, group_name)]
        lines.append(f"### {group_name} ({len(items)}个文件)\n")
        lines.append("| # | 操作 | 当前文件名 | 建议新名 | 备注 |")
        lines.append("|---|------|-----------|---------|------|")

        for idx, f, r in items:
            action = r.get("action", "skip")
            new_name = r.get("new_name", f["name"])
            reason = r.get("reason", "")
            # For archive action, show destination
            if action == "archive":
                new_name = "→ _历史版本/"
            elif action == "delete":
                new_name = "→ 删除"
            lines.append(f"| {idx + 1} | {action} | {f['name']} | {new_name} | {reason} |")

        lines.append("")

    plan_path.parent.mkdir(parents=True, exist_ok=True)
    plan_path.write_text("\n".join(lines), encoding="utf-8")
    logger.info("建议表已生成: %s", plan_path)


# ---------------------------------------------------------------------------
# Plan Execution
# ---------------------------------------------------------------------------
def parse_plan_md(plan_path: Path) -> list[dict]:
    """Parse the Markdown plan file and extract operations."""
    text = plan_path.read_text(encoding="utf-8")
    operations = []

    # Match table rows: | # | action | old_name | new_name | note |
    row_pattern = re.compile(
        r"\|\s*(\d+)\s*\|\s*(\w+)\s*\|\s*(.+?)\s*\|\s*(.+?)\s*\|\s*(.*?)\s*\|"
    )

    for line in text.splitlines():
        m = row_pattern.match(line.strip())
        if not m:
            continue
        idx_str, action, old_name, new_name, note = m.groups()
        # Skip header rows
        if action in ("操作", "---"):
            continue
        operations.append({
            "index": int(idx_str),
            "action": action.strip(),
            "old_name": old_name.strip(),
            "new_name": new_name.strip(),
            "note": note.strip(),
        })

    return operations


def execute_plan(
    files: list[dict],
    operations: list[dict],
    log_path: Path,
    logger: logging.Logger,
) -> None:
    """Execute the rename plan."""
    # Build index lookup
    file_map = {i + 1: f for i, f in enumerate(files)}

    changelog = []
    executed = 0
    skipped = 0

    for op in operations:
        idx = op["index"]
        action = op["action"]
        f = file_map.get(idx)

        if not f:
            logger.warning("索引 %d 不在文件列表中，跳过", idx)
            skipped += 1
            continue

        src = Path(f["path"])
        if not src.exists():
            logger.warning("文件不存在: %s，跳过", src)
            skipped += 1
            continue

        if action == "skip":
            skipped += 1
            continue

        elif action == "rename":
            new_name = op["new_name"]
            if new_name.startswith("→") or new_name == "—":
                skipped += 1
                continue
            dest = src.parent / new_name
            if dest.exists() and dest != src:
                logger.warning("目标已存在: %s，跳过", dest)
                skipped += 1
                continue
            if src.name == new_name:
                skipped += 1
                continue
            src.rename(dest)
            changelog.append({"action": "rename", "from": str(src), "to": str(dest)})
            logger.info("[重命名] %s → %s", src.name, new_name)
            executed += 1

        elif action == "archive":
            archive_dir = src.parent / "_历史版本"
            archive_dir.mkdir(exist_ok=True)
            dest = archive_dir / src.name
            if dest.exists():
                stem, ext = src.stem, src.suffix
                counter = 1
                while dest.exists():
                    dest = archive_dir / f"{stem} ({counter}){ext}"
                    counter += 1
            shutil.move(str(src), str(dest))
            changelog.append({"action": "archive", "from": str(src), "to": str(dest)})
            logger.info("[归档] %s → _历史版本/", src.name)
            executed += 1

        elif action == "delete":
            trash_dir = Path.home() / ".Trash"
            dest = trash_dir / src.name
            if dest.exists():
                stem, ext = src.stem, src.suffix
                counter = 1
                while dest.exists():
                    dest = trash_dir / f"{stem} ({counter}){ext}"
                    counter += 1
            shutil.move(str(src), str(dest))
            changelog.append({"action": "delete", "from": str(src), "to": str(dest)})
            logger.info("[删除] %s → ~/.Trash/", src.name)
            executed += 1

    # Save changelog for rollback
    log_path.write_text(json.dumps(changelog, ensure_ascii=False, indent=2), encoding="utf-8")
    logger.info("=== 执行完成: %d 个操作, %d 个跳过 ===", executed, skipped)
    logger.info("操作日志: %s", log_path)


def rollback(log_path: Path, logger: logging.Logger) -> None:
    """Rollback operations using the changelog."""
    if not log_path.exists():
        logger.error("操作日志不存在: %s", log_path)
        return

    changelog = json.loads(log_path.read_text(encoding="utf-8"))
    rolled = 0

    for entry in reversed(changelog):
        src = Path(entry["to"])
        dest = Path(entry["from"])
        if not src.exists():
            logger.warning("源文件不存在: %s，跳过回滚", src)
            continue
        dest.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(src), str(dest))
        logger.info("[回滚] %s → %s", src, dest)
        rolled += 1

    logger.info("=== 回滚完成: %d 个操作 ===", rolled)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def cmd_analyze(args, config, logger):
    """Analyze files and generate rename plan."""
    if args.dir:
        scan_dirs = [Path(args.dir).expanduser()]
    else:
        scan_dirs = [Path(d).expanduser() for d in config["scan_dirs"]]

    dry_run_limit = 5 if args.dry_run else 0

    logger.info("=== 开始分析 %s ===", "（dry-run 前5个文件）" if args.dry_run else "")

    files = scan_files(
        scan_dirs,
        config["ignore_prefixes"],
        config["ai"]["content_preview_chars"],
        logger,
        dry_run_limit,
    )

    if not files:
        logger.info("没有找到文件")
        return

    logger.info("共扫描 %d 个文件", len(files))

    # Find duplicates and similar names (for logging)
    dups = find_duplicates(files)
    if dups:
        logger.info("发现 %d 组 MD5 重复文件", len(dups))
    similar = find_similar_names(files, config["similarity_threshold"])
    if similar:
        logger.info("发现 %d 组文件名相似的文件", len(similar))

    # AI analysis
    ai_results = analyze_with_ai(files, config, logger)

    # Generate plan
    plan_path = Path(config["plan_file"]).expanduser()
    generate_plan_md(files, ai_results, plan_path, logger)

    # Also save raw file list for execute phase
    raw_path = plan_path.with_suffix(".json")
    raw_path.write_text(
        json.dumps(files, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    logger.info("文件元信息已保存: %s", raw_path)
    logger.info("请审核 %s 后运行: python3 smart_rename.py execute", plan_path)


def cmd_execute(args, config, logger):
    """Execute the approved rename plan."""
    plan_path = Path(args.plan or config["plan_file"]).expanduser()
    raw_path = plan_path.with_suffix(".json")
    log_path = Path(config["log_file"]).expanduser()

    if not plan_path.exists():
        logger.error("建议表不存在: %s", plan_path)
        logger.error("请先运行: python3 smart_rename.py analyze --all")
        return

    if not raw_path.exists():
        logger.error("文件元信息不存在: %s", raw_path)
        return

    files = json.loads(raw_path.read_text(encoding="utf-8"))
    operations = parse_plan_md(plan_path)

    if not operations:
        logger.info("建议表中没有操作")
        return

    logger.info("=== 开始执行: %d 个操作 ===", len(operations))
    execute_plan(files, operations, log_path, logger)


def cmd_rollback(args, config, logger):
    """Rollback the last execution."""
    log_path = Path(config["log_file"]).expanduser()
    rollback(log_path, logger)


def main():
    parser = argparse.ArgumentParser(description="Downloads 智能重命名工具")
    sub = parser.add_subparsers(dest="command", required=True)

    # analyze
    p_analyze = sub.add_parser("analyze", help="分析文件并生成重命名建议")
    p_analyze.add_argument("--all", action="store_true", help="扫描所有分类目录")
    p_analyze.add_argument("--dir", help="扫描指定目录")
    p_analyze.add_argument("--dry-run", action="store_true", help="只分析前5个文件（测试）")

    # execute
    p_execute = sub.add_parser("execute", help="执行审核后的重命名计划")
    p_execute.add_argument("--plan", help="指定建议表路径（默认 ~/Downloads/_rename_plan.md）")

    # rollback
    sub.add_parser("rollback", help="回滚上次执行")

    args = parser.parse_args()
    config = load_config()
    logger = setup_logging(config["app_log"])

    if args.command == "analyze":
        if not args.all and not args.dir:
            parser.error("请指定 --all 或 --dir")
        cmd_analyze(args, config, logger)
    elif args.command == "execute":
        cmd_execute(args, config, logger)
    elif args.command == "rollback":
        cmd_rollback(args, config, logger)


if __name__ == "__main__":
    main()
