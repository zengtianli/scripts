#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT文档字体统一工具 - 微软雅黑

功能：
- 将PPT文档中所有文字的字体设置为微软雅黑 (Microsoft YaHei)
- 处理所有幻灯片中的文本框、标题、占位符等
- 处理表格中的文字
- 保持所有其他格式属性（粗体、斜体、颜色、大小等）
- 自动备份原文件

用法:
    python3 format_font_yahei.py <input.pptx>

示例:
    python3 format_font_yahei.py presentation.pptx
"""

import sys
import os
import shutil
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent.parent / "lib"))
from finder import get_input_files

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    print("❌ 错误: 缺少 python-pptx 库")
    print("💡 请运行: pip install python-pptx")
    sys.exit(1)


# 目标字体
TARGET_FONT = 'Microsoft YaHei'

# XML 命名空间
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
}


def show_message(msg_type, message):
    """显示格式化消息"""
    icons = {
        'success': '✅',
        'error': '❌',
        'warning': '⚠️',
        'info': 'ℹ️',
        'processing': '🔄'
    }
    icon = icons.get(msg_type, 'ℹ️')
    print(f"{icon} {message}")


def backup_file(file_path):
    """备份原始文件"""
    backup_path = f"{file_path}.backup"
    try:
        shutil.copy2(file_path, backup_path)
        show_message('info', f"已备份原文件: {backup_path}")
        return backup_path
    except Exception as e:
        show_message('warning', f"备份文件失败: {e}")
        return None


def set_font_for_run(run, font_name=TARGET_FONT):
    """
    为run设置字体（强制XML级别设置）
    
    Args:
        run: pptx run对象
        font_name: 字体名称
    """
    try:
        # 1. 使用API设置
        run.font.name = font_name
        
        # 2. 强制XML级别设置 - 直接操作rPr元素
        rPr = run._r.get_or_add_rPr()
        
        # 查找或创建 a:latin 元素（西文字体）
        latin = rPr.find(qn('a:latin'))
        if latin is None:
            latin = etree.SubElement(rPr, qn('a:latin'))
        latin.set('typeface', font_name)
        
        # 查找或创建 a:ea 元素（东亚字体 - 中文）
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', font_name)
        
        # 查找或创建 a:cs 元素（复杂脚本字体）
        cs = rPr.find(qn('a:cs'))
        if cs is None:
            cs = etree.SubElement(rPr, qn('a:cs'))
        cs.set('typeface', font_name)
        
    except Exception as e:
        # 某些run可能没有字体属性
        pass


def set_font_for_paragraph_default(paragraph, font_name=TARGET_FONT):
    """
    设置段落的默认字体属性（defRPr）
    
    Args:
        paragraph: pptx paragraph对象
        font_name: 字体名称
    """
    try:
        pPr = paragraph._p.get_or_add_pPr()
        
        # 查找或创建 defRPr（默认文本属性）
        defRPr = pPr.find(qn('a:defRPr'))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn('a:defRPr'))
        
        # 设置 latin 字体
        latin = defRPr.find(qn('a:latin'))
        if latin is None:
            latin = etree.SubElement(defRPr, qn('a:latin'))
        latin.set('typeface', font_name)
        
        # 设置 ea 字体（东亚）
        ea = defRPr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(defRPr, qn('a:ea'))
        ea.set('typeface', font_name)
        
        # 设置 cs 字体
        cs = defRPr.find(qn('a:cs'))
        if cs is None:
            cs = etree.SubElement(defRPr, qn('a:cs'))
        cs.set('typeface', font_name)
        
    except Exception:
        pass


def set_font_for_endParaRPr(paragraph, font_name=TARGET_FONT):
    """
    设置段落结束符的字体属性（endParaRPr）
    
    Args:
        paragraph: pptx paragraph对象
        font_name: 字体名称
    """
    try:
        endParaRPr = paragraph._p.find(qn('a:endParaRPr'))
        if endParaRPr is not None:
            # 设置 latin 字体
            latin = endParaRPr.find(qn('a:latin'))
            if latin is None:
                latin = etree.SubElement(endParaRPr, qn('a:latin'))
            latin.set('typeface', font_name)
            
            # 设置 ea 字体
            ea = endParaRPr.find(qn('a:ea'))
            if ea is None:
                ea = etree.SubElement(endParaRPr, qn('a:ea'))
            ea.set('typeface', font_name)
            
            # 设置 cs 字体
            cs = endParaRPr.find(qn('a:cs'))
            if cs is None:
                cs = etree.SubElement(endParaRPr, qn('a:cs'))
            cs.set('typeface', font_name)
    except Exception:
        pass


def process_text_frame(text_frame, stats):
    """
    处理文本框中的所有段落和run
    
    Args:
        text_frame: pptx text_frame对象
        stats: 统计字典
    """
    for paragraph in text_frame.paragraphs:
        # 设置段落默认字体
        set_font_for_paragraph_default(paragraph, TARGET_FONT)
        # 设置段落结束符字体
        set_font_for_endParaRPr(paragraph, TARGET_FONT)
        
        # 处理每个run
        for run in paragraph.runs:
            set_font_for_run(run, TARGET_FONT)
            stats['processed_runs'] += 1


def process_table(table, stats):
    """
    处理表格中的所有单元格
    
    Args:
        table: pptx table对象
        stats: 统计字典
    """
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                process_text_frame(cell.text_frame, stats)
                stats['processed_tables'] += 1


def process_shape(shape, stats):
    """
    处理单个形状
    
    Args:
        shape: pptx shape对象
        stats: 统计字典
    """
    # 处理有文本框的形状
    if shape.has_text_frame:
        process_text_frame(shape.text_frame, stats)
        stats['processed_shapes'] += 1
    
    # 处理表格
    if shape.has_table:
        process_table(shape.table, stats)
    
    # 处理组合形状中的子形状
    if hasattr(shape, 'shapes'):
        for sub_shape in shape.shapes:
            process_shape(sub_shape, stats)


def process_slide(slide, stats):
    """
    处理单个幻灯片
    
    Args:
        slide: pptx slide对象
        stats: 统计字典
    """
    for shape in slide.shapes:
        process_shape(shape, stats)


def process_slide_master(slide_master, stats):
    """
    处理幻灯片母版
    
    Args:
        slide_master: pptx slide_master对象
        stats: 统计字典
    """
    # 处理母版中的形状
    for shape in slide_master.shapes:
        process_shape(shape, stats)
    
    # 处理母版的布局
    for layout in slide_master.slide_layouts:
        for shape in layout.shapes:
            process_shape(shape, stats)


def format_presentation(input_path, output_path=None):
    """
    格式化PPT文档中所有文字的字体
    
    Args:
        input_path: 输入文件路径
        output_path: 输出文件路径，如果为None则覆盖原文件
    
    Returns:
        bool: 是否成功
    """
    try:
        # 验证输入文件
        if not os.path.exists(input_path):
            show_message('error', f"文件不存在: {input_path}")
            return False
        
        if not input_path.lower().endswith('.pptx'):
            show_message('error', "只支持.pptx格式的文件")
            return False
        
        show_message('processing', f"正在处理文件: {os.path.basename(input_path)}")
        
        # 备份原文件
        backup_path = backup_file(input_path)
        
        # 打开PPT
        prs = Presentation(input_path)
        
        # 统计信息
        total_slides = len(prs.slides)
        stats = {
            'processed_shapes': 0,
            'processed_runs': 0,
            'processed_tables': 0
        }
        
        show_message('info', f"文档包含 {total_slides} 张幻灯片")
        
        # 处理幻灯片母版（重要：这里的字体设置会影响整个PPT）
        show_message('processing', "正在处理幻灯片母版...")
        for slide_master in prs.slide_masters:
            try:
                process_slide_master(slide_master, stats)
            except Exception as e:
                show_message('warning', f"处理母版时出错: {e}")
        
        # 处理所有幻灯片
        show_message('processing', "正在处理幻灯片...")
        for i, slide in enumerate(prs.slides, 1):
            try:
                process_slide(slide, stats)
            except Exception as e:
                show_message('warning', f"处理第{i}张幻灯片时出错: {e}")
                continue
        
        show_message('info', f"已处理 {stats['processed_shapes']} 个形状, {stats['processed_runs']} 个文本run")
        if stats['processed_tables'] > 0:
            show_message('info', f"已处理 {stats['processed_tables']} 个表格单元格")
        
        # 保存文档
        output_file = output_path if output_path else input_path
        prs.save(output_file)
        
        show_message('success', f"字体格式化完成: {os.path.basename(output_file)}")
        show_message('info', f"所有文字已设置为: {TARGET_FONT}")
        
        if backup_path:
            show_message('info', f"如需恢复，请使用备份文件: {os.path.basename(backup_path)}")
        
        return True
        
    except Exception as e:
        show_message('error', f"处理文件时出错: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """主函数"""
    # 获取输入文件（优先命令行参数，否则从 Finder 获取）
    files = get_input_files(sys.argv[1:], expected_ext='pptx', allow_multiple=False)
    
    if not files:
        print("PPT文档字体统一工具 - 微软雅黑")
        print()
        print("用法:")
        print(f"  python {os.path.basename(__file__)} <输入文件.pptx>")
        print(f"  python {os.path.basename(__file__)} <输入文件.pptx> <输出文件.pptx>")
        print("  或在 Finder 中选择 .pptx 文件后运行")
        print()
        print("功能:")
        print(f"  将PPT文档中所有文字的字体设置为 {TARGET_FONT}")
        print("  保持文字的其他格式属性（粗体、斜体、颜色、大小等）")
        print()
        print("示例:")
        print(f"  python {os.path.basename(__file__)} presentation.pptx")
        print(f"  python {os.path.basename(__file__)} input.pptx output.pptx")
        sys.exit(1)
    
    input_path = files[0]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None
    
    # 处理文件
    success = format_presentation(input_path, output_path)
    
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()

